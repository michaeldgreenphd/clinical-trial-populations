"""Fill slides 1 and 2 of the Anthropic Healthcare deck.

    python3 slides/fill_user_deck.py <source.pptx> [-o out.pptx]

Those two slides ship as title-only placeholders. This adds their content on the
same grid the rest of the deck uses (measured off slide 6):

    title      18pt bold      y=0.29     margins x=0.35 .. 9.65
    subtitle   10.5pt         y=0.69
    col header 11pt           y=1.02     hairline rule beneath at y=1.32
    body       8.5pt / 9pt
    footnote   7pt            y=5.26
    columns    split at x=5.90 with a vertical hairline

Every figure is sourced in a comment beside it. Canvas is 10 x 5.625in.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

IMG = Path(__file__).parent / "img"

FONT = "Century Schoolbook"
INK = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x6D, 0x6C, 0x66)
GREEN = RGBColor(0x15, 0x80, 0x3D)
TERRA = RGBColor(0xB4, 0x48, 0x3D)
RULE = RGBColor(0xD9, 0xD8, 0xD3)


def textbox(slide, x, y, w, h, runs, *, size=8.5, color=INK, bold=False,
            italic=False, space_after=0, line_spacing=None, align=None):
    """Add a text box. `runs` is a string or a list of (text, overrides) pairs;
    a run of "\n" starts a new paragraph."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [(runs, {})]

    para = tf.paragraphs[0]
    if align:
        para.alignment = align
    if line_spacing:
        para.line_spacing = line_spacing
    for text, over in runs:
        if text == "\n":
            para = tf.add_paragraph()
            if align:
                para.alignment = align
            if line_spacing:
                para.line_spacing = line_spacing
            para.space_before = Pt(space_after)
            continue
        run = para.add_run()
        run.text = text
        f = run.font
        f.name = FONT
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.italic = over.get("italic", italic)
        f.color.rgb = over.get("color", color)
    return box


def hairline(slide, x, y, w, color=RULE, weight=0.9):
    """Thin horizontal rule — the divider idiom used on slide 6."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Pt(weight))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def vrule(slide, x, y, h, color=RULE, weight=0.9):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                 Pt(weight), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def restyle_title(slide, text):
    """Normalise the placeholder title onto the slide-6 grid."""
    box = slide.shapes[0]
    box.left, box.top = Inches(0.35), Inches(0.29)
    box.width, box.height = Inches(9.30), Inches(0.42)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    para = tf.paragraphs[0]
    for run in list(para.runs)[1:]:
        run._r.getparent().remove(run._r)
    run = para.runs[0] if para.runs else para.add_run()
    run.text = text
    run.font.name, run.font.size, run.font.bold = FONT, Pt(18), True
    run.font.color.rgb = INK


# --------------------------------------------------------------------------- #
# Slide 1 — what the clinical-trial work established, and the pivot to AI
# --------------------------------------------------------------------------- #
def build_slide1(slide):
    restyle_title(slide, "Conclusions from Civic Sample Project so Far")
    textbox(slide, 0.35, 0.69, 9.30, 0.28, [
        ("Four demographic dimensions across ", {}),
        ("79,297 trials with posted results", {"bold": True}),
        (" — and the same questions now aimed at FDA-authorized medical AI.", {}),
    ], size=10.5, color=MUTED, italic=True)

    # ---- Band A: the four dashboard dimensions ----------------------------- #
    textbox(slide, 0.35, 1.00, 6.0, 0.26,
            "What we built for clinical trials, 2009–2026", size=11, color=INK)
    hairline(slide, 0.35, 1.32, 9.30)

    # The three doughnuts are captured with their slices labelled in place by
    # chartjs-plugin-datalabels (see capture_screenshots.js), so no separate key
    # is needed. Shares are the all-study-types cohort — capture_screenshots.js
    # switches the dashboard's Study Type filter to "all" before shooting, so the
    # on-chart percentages match the 79,297-trial headline.
    dims = [
        {"img": "crop-race.png", "ratio": 912 / 671, "label": "RACE", "star": "",
         "stat": "57.0%", "of": "of trials report it"},
        {"img": "crop-ethnicity.png", "ratio": 911 / 564, "label": "ETHNICITY", "star": "",
         "stat": "39.1%", "of": "of trials report it"},
        {"img": "crop-sex.png", "ratio": 755 / 574, "label": "SEX", "star": "*",
         "stat": "97.3%", "of": "of trials report it"},
        {"img": "crop-geography.png", "ratio": 1850 / 1150, "label": "GEOGRAPHY", "star": "*",
         "stat": "37.3%", "of": "of sites are in the South · Midwest 23.6% · "
                              "West 20.6% · Northeast 18.5%"},
    ]

    col_w, band_h, band_y = 2.20, 1.52, 1.42
    gap = (9.30 - 4 * col_w) / 3
    for i, d in enumerate(dims):
        x = 0.35 + i * (col_w + gap)
        w, h = band_h * d["ratio"], band_h
        if w > col_w:                       # wide capture — size by column instead
            w, h = col_w, col_w / d["ratio"]
        slide.shapes.add_picture(str(IMG / d["img"]),
                                 Inches(x + (col_w - w) / 2),
                                 Inches(band_y + (band_h - h) / 2),
                                 width=Inches(w), height=Inches(h))

        textbox(slide, x, 3.02, col_w, 0.14, [
            (d["label"], {}),
            (d["star"], {"color": TERRA}),
        ], size=7, color=MUTED, bold=True)
        textbox(slide, x, 3.16, col_w, 0.40, [
            (d["stat"], {"size": 9.5, "bold": True, "color": GREEN}),
            ("  " + d["of"], {"size": 6.5, "color": MUTED}),
        ], line_spacing=0.95)

    # ---- Band B: the pivot to medical AI ----------------------------------- #
    hairline(slide, 0.35, 3.60, 9.30)
    textbox(slide, 0.35, 3.68, 6.0, 0.26,
            "The same four questions, now for medical AI", size=11, color=INK)

    # Source: ai-ml-enabled-devices-csv_20260305.csv (FDA list, 2026-03-05).
    slide.shapes.add_picture(str(IMG / "crop-ai-devices.png"), Inches(0.35),
                             Inches(4.00), width=Inches(4.55), height=Inches(0.87))

    vrule(slide, 5.20, 3.98, 1.02)

    textbox(slide, 5.42, 3.98, 4.23, 1.10, [
        ("96.2% arrive by the 510(k) predicate route", {"size": 9, "bold": True, "color": TERRA}),
        (" — substantial equivalence to a device already marketed, with no new "
         "clinical evidence required. Only 18 of 1,451 went through full PMA review.", {}),
        ("\n", {}),
        ("The FDA publishes no demographics for these validation cohorts.", {"size": 9, "bold": True}),
        (" So Claude reads each device’s public 510(k)/De Novo/PMA summary and its "
         "open-access manuscripts, quoting the evidence before committing a value.", {}),
    ], size=8, color=INK, space_after=4, line_spacing=0.95)

    textbox(slide, 0.35, 5.12, 9.30, 0.14, [
        ("*", {"color": TERRA, "bold": True}),
        (" Sex and geography sit outside the race/ethnicity analyses already written up — "
         "protocol manuscripts covering their extraction are in preparation.", {}),
    ], size=6.8, color=MUTED)
    textbox(slide, 0.35, 5.28, 9.30, 0.22,
            "Dashboard: civicsample.com, all study types (79,297 trials) · slices under 3% "
            "are left unlabelled · geography shares are US Census regions "
            "· device figures from the FDA AI/ML-Enabled Device List",
            size=6.8, color=MUTED)


# --------------------------------------------------------------------------- #
# Slide 2 — token usage against the award, and the open questions
# --------------------------------------------------------------------------- #
def placeholder(slide, x, y, w, h, lines):
    """Dashed box marking content the team still has to supply."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFA, 0xF7, 0xF6)
    box.line.color.rgb = TERRA
    box.line.width = Pt(0.75)
    try:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        box.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    except Exception:
        pass                      # solid border is an acceptable fallback
    box.shadow.inherit = False
    textbox(slide, x + 0.12, y + 0.10, w - 0.24, h - 0.20, lines,
            size=7.5, color=MUTED, space_after=3, line_spacing=0.98)
    return box


def build_slide2(slide):
    restyle_title(slide, "What we did with tokens, and where feedback would be most useful")
    textbox(slide, 0.35, 0.69, 9.30, 0.28, [
        ("What it costs to pull one concept out of a regulatory PDF", {"bold": True}),
        (" — and where feedback would help most.", {}),
    ], size=10.5, color=MUTED, italic=True)

    # ---- Band 1, left: the prompt and what it returns ----------------------- #
    textbox(slide, 0.35, 1.00, 5.15, 0.26, "The extraction, end to end", size=11, color=INK)
    hairline(slide, 0.35, 1.32, 5.15)

    textbox(slide, 0.35, 1.40, 5.15, 0.13, "1 · THE REVIEW PROMPT", size=7,
            color=MUTED, bold=True)
    placeholder(slide, 0.35, 1.54, 5.15, 1.02, [
        ("PLACEHOLDER — Maryam", {"bold": True, "color": TERRA}),
        ("\n", {}),
        ("Paste the review prompt here: the system prompt that tells Claude to quote "
         "its evidence before committing a value, plus whatever review criteria you "
         "want to show the room.", {}),
    ])

    textbox(slide, 0.35, 2.64, 5.15, 0.13, "2 · WHAT COMES BACK", size=7,
            color=MUTED, bold=True)

    # Real Sonnet 4.6 output for DEN140025, from data/fda_demographics_extracted.json.
    out = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(2.78),
                                 Inches(5.15), Inches(0.70))
    out.fill.solid()
    out.fill.fore_color.rgb = RGBColor(0xF6, 0xF7, 0xF5)
    out.line.color.rgb = RULE
    out.line.width = Pt(0.75)
    out.shadow.inherit = False
    textbox(slide, 0.47, 2.85, 4.91, 0.60, [
        ("DEN140025 · BrainScope Ahead 100 · 6-page De Novo summary", {"bold": True, "size": 7}),
        ("\n", {}),
        ("device_name  ", {"size": 6.8, "color": MUTED}),
        ("“BrainScope Ahead 100 (Ahead® M-100 and CV-100)”", {"size": 6.8}),
        ("  p.1", {"size": 6.8, "color": MUTED}),
        ("\n", {}),
        ("age_range    ", {"size": 6.8, "color": MUTED}),
        ("“18–80 years”", {"size": 6.8}),
        ("  ← quoted from “…are between the ages of 18-80 years.”  p.1", {"size": 6.8, "color": MUTED}),
        ("\n", {}),
        ("sex · race · ethnicity  ", {"size": 6.8, "color": MUTED}),
        ("not reported anywhere in the document", {"size": 6.8, "bold": True, "color": TERRA}),
    ], size=6.8, color=INK, space_after=2, line_spacing=0.98)

    # ---- Band 1, right: what a concept costs -------------------------------- #
    vrule(slide, 5.62, 1.40, 2.08)
    textbox(slide, 5.80, 1.00, 3.85, 0.26, "What a concept costs to extract",
            size=11, color=INK)
    hairline(slide, 5.80, 1.32, 3.85)

    cols = [(5.80, 1.90, PP_ALIGN.LEFT), (7.72, 0.60, PP_ALIGN.RIGHT),
            (8.36, 0.60, PP_ALIGN.RIGHT), (9.02, 0.63, PP_ALIGN.RIGHT)]
    heads = ["CONCEPT GROUP", "HAIKU", "SONNET", "OPUS"]
    for (cx, cw, al), head in zip(cols, heads):
        textbox(slide, cx, 1.42, cw, 0.13, head, size=6, color=MUTED, bold=True, align=al)
    hairline(slide, 5.80, 1.57, 3.85)

    # Real figures: data/fda_token_metrics.json at current list rates, per 1,000 docs.
    rows = [("Current FDA schema — 20 fields", "$23.95", "$62.99", "$141.70", False),
            ("[ Maryam — concept group ]", "", "", "", True),
            ("[ Maryam — concept group ]", "", "", "", True),
            ("[ Maryam — concept group ]", "", "", "", True)]
    ry = 1.63
    for name, h, s, o, is_ph in rows:
        colour = TERRA if is_ph else INK
        textbox(slide, cols[0][0], ry, cols[0][1], 0.13, name, size=6.8,
                color=colour, italic=is_ph)
        for (cx, cw, al), val in zip(cols[1:], (h, s, o)):
            textbox(slide, cx, ry, cw, 0.13, val or "—", size=6.8,
                    color=INK if val else RULE, bold=bool(val), align=al)
        ry += 0.165
    hairline(slide, 5.80, ry + 0.02, 3.85)

    textbox(slide, 5.80, ry + 0.10, 3.85, 0.30, [
        ("Cost per 1,000 documents at current list rates. ", {}),
        ("Amortized over the 20 fields the schema returns, that is $0.0012 per field "
         "on Haiku and $0.0071 on Opus.", {"bold": True, "color": INK}),
    ], size=6.5, color=MUTED, line_spacing=0.98)

    placeholder(slide, 5.80, 2.86, 3.85, 0.62, [
        ("PLACEHOLDER — efficiency", {"bold": True, "color": TERRA}),
        ("\n", {}),
        ("Manual review takes ____ per document; the pipeline takes ____. "
         "Add the human-hours baseline you want to compare against.", {}),
    ])

    # ---- Band 2: what we need a read on, and what this offers back ---------- #
    hairline(slide, 0.35, 3.56, 9.30)
    textbox(slide, 0.35, 3.62, 7.5, 0.26,
            "What we’d want your read on — and where this could be useful to you",
            size=11, color=INK)

    # Three questions we cannot settle from the data alone, then the offer.
    # Each carries the question, the detail, and the stake — what turns on the answer.
    asks = [
        ("CLINICAL", TERRA, "Representative against what?",
         "We can report that a validation cohort was 84% White. We cannot say whether "
         "that is a problem. The census, the condition’s own epidemiology and the "
         "deploying hospital’s panel all give different answers.",
         "Until this is settled we can publish composition, but not a judgment about it."),
        ("CLINICAL", TERRA, "Which social factors change a deployment decision?",
         "The schema reserves fields for income, education, insurance and neighborhood, "
         "but each one costs tokens and curator review to keep.",
         "Your answer decides what the rest of the schema budget buys."),
        ("TECHNICAL", TERRA, "Silence vs. a missed read.",
         "Quoting evidence before committing a value makes a positive claim auditable. "
         "But a document that is genuinely silent and a page the model skipped look "
         "identical in the output. How would you prompt or verify to tell them apart?",
         "Get this wrong and the central claim — that demographics are missing — is "
         "unsafe to make."),
        ("FOR ANTHROPIC", GREEN, "What this could give your team.",
         "Health systems ask whether a tool was validated on patients like theirs. The "
         "public record cannot answer that today. We are building the answer for all "
         "1,451 authorized devices — independent of the vendor, each value traceable to "
         "its source page.",
         "If that is evidence your patient-safety or policy conversations need, tell us "
         "and we will shape the first release around it."),
    ]
    qw, qgap = 2.20, (9.30 - 4 * 2.20) / 3
    for i, (tag, tag_col, q, sub, stake) in enumerate(asks, 1):
        qx = 0.35 + (i - 1) * (qw + qgap)
        textbox(slide, qx, 3.88, qw, 0.12, tag, size=6, bold=True, color=tag_col)
        textbox(slide, qx, 4.00, 0.16, 0.13, str(i), size=7.5, bold=True, color=tag_col)
        textbox(slide, qx + 0.16, 4.00, qw - 0.16, 0.26, q, size=7.5, bold=True,
                color=INK, line_spacing=0.95)
        textbox(slide, qx + 0.16, 4.28, qw - 0.16, 0.62, sub, size=6.4, color=MUTED,
                line_spacing=0.98)
        textbox(slide, qx + 0.16, 4.94, qw - 0.16, 0.28, stake, size=6.4, color=INK,
                italic=True, line_spacing=0.98)

    textbox(slide, 0.35, 5.30, 9.30, 0.22,
            "Pilot to date: 36 documents · 373 pages · 2.63M tokens, each document run "
            "through all three models · token counts from data/fda_token_metrics.json "
            "· costs at Anthropic list rates · extraction sample is real Sonnet 4.6 output",
            size=6.8, color=MUTED)


def main():
    src = Path(sys.argv[1])
    out = Path(sys.argv[sys.argv.index("-o") + 1]) if "-o" in sys.argv \
        else src.with_name("Presentation_for_Anthropic_Healthcare_filled.pptx")
    prs = Presentation(str(src))
    build_slide1(prs.slides[0])
    build_slide2(prs.slides[1])
    prs.save(str(out))
    print(f"wrote {out}")


if __name__ == "__main__":
    main()
